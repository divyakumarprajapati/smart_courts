#!/usr/bin/env python3
"""
CourtNG Pitch Deck Generator
Creates a professional PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn, nsmap
import colorsys

# Brand Colors
PRIMARY = RgbColor(16, 185, 129)      # #10b981 - Emerald
PRIMARY_DARK = RgbColor(5, 150, 105)   # #059669
PRIMARY_LIGHT = RgbColor(52, 211, 153) # #34d399
ACCENT = RgbColor(99, 102, 241)        # #6366f1 - Indigo
DARK = RgbColor(15, 23, 42)            # #0f172a
DARK_SECONDARY = RgbColor(30, 41, 59)  # #1e293b
GRAY_400 = RgbColor(156, 163, 175)     # #9ca3af
GRAY_500 = RgbColor(107, 114, 128)     # #6b7280
GRAY_600 = RgbColor(75, 85, 99)        # #4b5563
WHITE = RgbColor(255, 255, 255)
RED = RgbColor(239, 68, 68)            # #ef4444
RED_LIGHT = RgbColor(252, 165, 165)    # #fca5a5


def set_shape_gradient(shape, color1, color2, angle=135):
    """Set gradient fill on a shape"""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2


def set_shape_solid(shape, color):
    """Set solid fill on a shape"""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_frame(shape, text, font_size=18, font_color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    """Add formatted text to a shape"""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = "Arial"
    return tf


def add_paragraph(text_frame, text, font_size=14, font_color=WHITE, bold=False, space_before=0):
    """Add a new paragraph to text frame"""
    p = text_frame.add_paragraph()
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = "Arial"
    return p


def create_card(slide, left, top, width, height, fill_color, border_color=None):
    """Create a rounded rectangle card"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    set_shape_solid(shape, fill_color)
    shape.line.fill.background()
    
    # Adjust corner radius
    shape.adjustments[0] = 0.1
    return shape


def create_slide_background(slide, color1, color2):
    """Set gradient background for slide"""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 135
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2


def add_slide_header(slide, slide_num, total=12):
    """Add consistent header to slides"""
    # Slide number
    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(1.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{slide_num:02d} / {total:02d}"
    run.font.size = Pt(11)
    run.font.color.rgb = GRAY_500
    run.font.bold = True
    run.font.name = "Arial"
    
    # Logo text on right
    logo_box = slide.shapes.add_textbox(Inches(8.2), Inches(0.4), Inches(1.5), Inches(0.3))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "CourtNG"
    run.font.size = Pt(14)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Arial"


def add_slide_title(slide, title, subtitle=None, accent_word=None):
    """Add title and optional subtitle"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    
    if accent_word and accent_word in title:
        parts = title.split(accent_word)
        run = p.add_run()
        run.text = parts[0]
        run.font.size = Pt(40)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Arial"
        
        run2 = p.add_run()
        run2.text = accent_word
        run2.font.size = Pt(40)
        run2.font.color.rgb = PRIMARY_LIGHT
        run2.font.bold = True
        run2.font.name = "Arial"
        
        if len(parts) > 1:
            run3 = p.add_run()
            run3.text = parts[1]
            run3.font.size = Pt(40)
            run3.font.color.rgb = WHITE
            run3.font.bold = True
            run3.font.name = "Arial"
    else:
        run = p.add_run()
        run.text = title
        run.font.size = Pt(40)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Arial"
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(8), Inches(0.6))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(16)
        run.font.color.rgb = GRAY_400
        run.font.name = "Arial"


# =============================================================================
# SLIDE CREATION FUNCTIONS
# =============================================================================

def create_title_slide(prs):
    """Slide 1: Title/Start Slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    create_slide_background(slide, DARK, DARK_SECONDARY)
    
    # Logo circle
    logo = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.15), Inches(1.2), Inches(1.5), Inches(1.5)
    )
    set_shape_gradient(logo, PRIMARY, RgbColor(6, 182, 212))
    logo.line.fill.background()
    logo.adjustments[0] = 0.3
    
    # Court icon text (simplified)
    icon_text = slide.shapes.add_textbox(Inches(4.15), Inches(1.5), Inches(1.5), Inches(0.9))
    tf = icon_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "⊞"
    run.font.size = Pt(48)
    run.font.color.rgb = WHITE
    
    # Title: CourtNG
    title = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Court"
    run.font.size = Pt(72)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Arial"
    run2 = p.add_run()
    run2.text = "NG"
    run2.font.size = Pt(72)
    run2.font.color.rgb = PRIMARY_LIGHT
    run2.font.bold = True
    run2.font.name = "Arial"
    
    # Tagline
    tagline = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.5))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "The Operating System for Intelligent Courts"
    run.font.size = Pt(24)
    run.font.color.rgb = GRAY_400
    run.font.bold = True
    run.font.name = "Arial"
    
    # Description
    desc = slide.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(7), Inches(0.6))
    tf = desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "AI-powered court automation that officiates games, calls every line, and delivers real-time analytics — without human intervention."
    run.font.size = Pt(14)
    run.font.color.rgb = GRAY_500
    run.font.name = "Arial"
    
    # Badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(5.5), Inches(2.8), Inches(0.4)
    )
    set_shape_solid(badge, RgbColor(30, 41, 59))
    badge.line.color.rgb = GRAY_600
    badge.adjustments[0] = 0.5
    
    badge_text = slide.shapes.add_textbox(Inches(3.5), Inches(5.55), Inches(2.8), Inches(0.4))
    tf = badge_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "● Seed Stage • 2026"
    run.font.size = Pt(12)
    run.font.color.rgb = GRAY_400
    run.font.name = "Arial"


def create_problem_slide(prs):
    """Slide 2: Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_background(slide, DARK, DARK_SECONDARY)
    add_slide_header(slide, 2)
    add_slide_title(slide, "The Problem", 
                   "Sports courts today are stuck in the past — expensive, inconsistent, and data-blind.")
    
    problems = [
        ("⚠️", "Human Error & Disputes", "Manual scoring and refereeing leads to errors, disputes, and inconsistent rule enforcement."),
        ("💰", "Expensive Officials", "Qualified referees are scarce and expensive. Venues spend $172K/month on staff overhead."),
        ("📊", "Zero Player Insights", "No real-time data means missed coaching opportunities and no performance analytics."),
        ("🚫", "Unaffordable Technology", "Hawk-Eye costs $60K+ per court. Only elite venues can afford automation."),
    ]
    
    positions = [(0.5, 2.3), (5, 2.3), (0.5, 4.1), (5, 4.1)]
    
    for i, (icon, title, desc) in enumerate(problems):
        left, top = positions[i]
        
        # Card background
        card = create_card(slide, left, top, 4.3, 1.5, RgbColor(50, 30, 30))
        card.line.color.rgb = RgbColor(120, 50, 50)
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.3), Inches(0.5), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = icon
        run.font.size = Pt(24)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(left + 0.8), Inches(top + 0.25), Inches(3.3), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(16)
        run.font.color.rgb = RED_LIGHT
        run.font.bold = True
        run.font.name = "Arial"
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.8), Inches(top + 0.7), Inches(3.3), Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(12)
        run.font.color.rgb = GRAY_400
        run.font.name = "Arial"


def create_solution_slide(prs):
    """Slide 3: Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_background(slide, RgbColor(15, 35, 35), DARK_SECONDARY)
    add_slide_header(slide, 3)
    add_slide_title(slide, "The Solution: CourtNG", 
                   "An AI-first court operating system that fully automates game flow and officiating.",
                   accent_word="CourtNG")
    
    solutions = [
        ("🤖", "Autonomous Officiating", "AI detects rally start/end, scores points, and calls fouls in real-time — no humans needed."),
        ("🎯", "Sport-Agnostic Engine", "Configure once, deploy anywhere: tennis, pickleball, badminton, futsal, and more."),
        ("⚡", "Edge-First Deployment", "On-premise NVIDIA Jetson inference. Sub-50ms latency. Works offline."),
        ("📈", "Integrated Analytics", "Live player stats, heatmaps, and performance trends delivered in real-time."),
    ]
    
    positions = [(0.5, 2.3), (5, 2.3), (0.5, 4.1), (5, 4.1)]
    
    for i, (icon, title, desc) in enumerate(solutions):
        left, top = positions[i]
        
        # Card background
        card = create_card(slide, left, top, 4.3, 1.5, RgbColor(20, 50, 45))
        card.line.color.rgb = RgbColor(40, 100, 90)
        
        # Icon with gradient background
        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + 0.2), Inches(top + 0.3), Inches(0.5), Inches(0.5)
        )
        set_shape_gradient(icon_bg, PRIMARY, RgbColor(6, 182, 212))
        icon_bg.line.fill.background()
        icon_bg.adjustments[0] = 0.2
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.3), Inches(0.5), Inches(0.5))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = icon
        run.font.size = Pt(20)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(left + 0.85), Inches(top + 0.25), Inches(3.3), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(16)
        run.font.color.rgb = PRIMARY_LIGHT
        run.font.bold = True
        run.font.name = "Arial"
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.85), Inches(top + 0.7), Inches(3.3), Inches(0.7))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(12)
        run.font.color.rgb = GRAY_400
        run.font.name = "Arial"


def create_validation_slide(prs):
    """Slide 4: Market Validation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_background(slide, DARK, DARK_SECONDARY)
    add_slide_header(slide, 4)
    add_slide_title(slide, "Market Validation", 
                   "The elite tier has proven AI officiating works. Now it's time to democratize access.")
    
    validations = [
        ("🏆", "Elite Tier Proven", "Wimbledon (2025), ATP, and US Open have deployed full AI line-calling. The technology works at the highest level."),
        ("📈", "Mid-Market Growth", "PlaySight, Wingfield, SwingVision, and TVConal are all raising funding — proving market demand."),
        ("🎯", "The Gap We're Filling", "No vendor offers fully autonomous, sport-agnostic \"court OS\". Existing players focus on single-sport, analytics-only solutions."),
    ]
    
    for i, (icon, title, desc) in enumerate(validations):
        top = 2.3 + (i * 1.25)
        
        # Card background
        card = create_card(slide, 0.5, top, 8.8, 1.05, RgbColor(25, 35, 50))
        card.line.color.rgb = RgbColor(50, 70, 100)
        
        # Icon with gradient background
        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(top + 0.25), Inches(0.55), Inches(0.55)
        )
        set_shape_gradient(icon_bg, PRIMARY, RgbColor(6, 182, 212))
        icon_bg.line.fill.background()
        icon_bg.adjustments[0] = 0.2
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(0.7), Inches(top + 0.27), Inches(0.55), Inches(0.55))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = icon
        run.font.size = Pt(22)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1.4), Inches(top + 0.2), Inches(7.5), Inches(0.35))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Arial"
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.4), Inches(top + 0.55), Inches(7.5), Inches(0.45))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(12)
        run.font.color.rgb = GRAY_400
        run.font.name = "Arial"


def create_market_size_slide(prs):
    """Slide 5: Market Size"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_background(slide, RgbColor(20, 20, 50), DARK_SECONDARY)
    add_slide_header(slide, 5)
    add_slide_title(slide, "Market Size", 
                   "A massive, rapidly growing opportunity across multiple converging markets.")
    
    stats = [
        ("$8.7B", "Computer Vision in Sports", "Projected by 2029"),
        ("$10.1B", "Tournament Management Software", "2033 • 14.7% CAGR"),
        ("$991B", "Sports Facilities Global", "2034 • 24.3% CAGR"),
    ]
    
    for i, (value, label, source) in enumerate(stats):
        left = 0.5 + (i * 3.2)
        
        # Card
        card = create_card(slide, left, 2.3, 2.9, 2, RgbColor(30, 40, 60))
        card.line.color.rgb = RgbColor(60, 80, 120)
        
        # Value
        val_box = slide.shapes.add_textbox(Inches(left), Inches(2.5), Inches(2.9), Inches(0.7))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = value
        run.font.size = Pt(40)
        run.font.color.rgb = PRIMARY_LIGHT
        run.font.bold = True
        run.font.name = "Arial"
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(3.2), Inches(2.7), Inches(0.5))
        tf = label_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(13)
        run.font.color.rgb = GRAY_400
        run.font.name = "Arial"
        
        # Source
        source_box = slide.shapes.add_textbox(Inches(left), Inches(3.75), Inches(2.9), Inches(0.3))
        tf = source_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = source
        run.font.size = Pt(10)
        run.font.color.rgb = GRAY_600
        run.font.name = "Arial"
    
    # Bottom stats bar
    bar = create_card(slide, 0.5, 4.6, 8.8, 0.9, RgbColor(30, 40, 60))
    bar.line.color.rgb = RgbColor(60, 80, 120)
    
    bottom_stats = [("50K+", "Addressable Venues"), ("159%", "Pickleball Growth (India)"), ("20K+", "US Courts by 2026")]
    
    for i, (val, lbl) in enumerate(bottom_stats):
        left = 1 + (i * 2.9)
        
        val_box = slide.shapes.add_textbox(Inches(left), Inches(4.7), Inches(2.5), Inches(0.4))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = val
        run.font.size = Pt(24)
        run.font.color.rgb = PRIMARY_LIGHT
        run.font.bold = True
        run.font.name = "Arial"
        
        lbl_box = slide.shapes.add_textbox(Inches(left), Inches(5.1), Inches(2.5), Inches(0.3))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = lbl
        run.font.size = Pt(10)
        run.font.color.rgb = GRAY_400
        run.font.name = "Arial"


def create_product_features_slide(prs):
    """Slide 6: Product - Core Features"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_background(slide, DARK, DARK_SECONDARY)
    add_slide_header(slide, 6)
    add_slide_title(slide, "Product: Core Capabilities", 
                   "Everything a court needs to run itself — in one unified platform.")
    
    features = [
        ("👁️", "Computer Vision", "Multi-camera ball & player tracking with sub-centimeter precision."),
        ("⚙️", "Rules Engine", "Declarative, configurable rules for any sport."),
        ("🎮", "Event Detection", "Real-time rally detection, scoring, game flow."),
        ("📊", "Live Analytics", "Player stats, heatmaps, performance metrics."),
        ("🖥️", "Admin Dashboard", "Full control over courts, matches, and settings."),
        ("📱", "Mobile Apps", "Player apps for stats, booking, and social."),
    ]
    
    for i, (icon, title, desc) in enumerate(features):
        col = i % 3
        row = i // 3
        left = 0.5 + (col * 3.15)
        top = 2.3 + (row * 1.7)
        
        # Card
        card = create_card(slide, left, top, 2.95, 1.5, RgbColor(25, 35, 50))
        card.line.color.rgb = RgbColor(50, 70, 100)
        
        # Icon with gradient background
        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + 0.95), Inches(top + 0.15), Inches(0.65), Inches(0.65)
        )
        set_shape_gradient(icon_bg, PRIMARY, RgbColor(6, 182, 212))
        icon_bg.line.fill.background()
        icon_bg.adjustments[0] = 0.25
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(left + 0.95), Inches(top + 0.18), Inches(0.65), Inches(0.65))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = icon
        run.font.size = Pt(24)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.85), Inches(2.75), Inches(0.3))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(15)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Arial"
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 1.15), Inches(2.75), Inches(0.35))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(10)
        run.font.color.rgb = GRAY_400
        run.font.name = "Arial"


def create_architecture_slide(prs):
    """Slide 7: Product - Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_background(slide, RgbColor(15, 35, 35), DARK_SECONDARY)
    add_slide_header(slide, 7)
    add_slide_title(slide, "Product: Architecture", 
                   "Modular, upgradeable, and built for scale from day one.")
    
    # Layer 1: Top
    modules_top = [
        ("📹 Vision Module", "Multi-camera, ball/player tracking"),
        ("📐 Rules Engine", "Sport-agnostic, declarative rules"),
    ]
    
    for i, (title, desc) in enumerate(modules_top):
        left = 1 + (i * 4)
        card = create_card(slide, left, 2.3, 3.6, 1, RgbColor(25, 40, 50))
        card.line.color.rgb = RgbColor(50, 80, 100)
        
        title_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(2.4), Inches(3.4), Inches(0.35))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(14)
        run.font.color.rgb = PRIMARY_LIGHT
        run.font.bold = True
        run.font.name = "Arial"
        
        desc_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(2.8), Inches(3.4), Inches(0.35))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(11)
        run.font.color.rgb = GRAY_500
        run.font.name = "Arial"
    
    # Layer 2: Core (highlighted)
    core = create_card(slide, 1, 3.5, 7.6, 0.9, PRIMARY)
    core.line.fill.background()
    
    core_title = slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(7.6), Inches(0.35))
    tf = core_title.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "🧠 CourtNG Core"
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Arial"
    
    core_desc = slide.shapes.add_textbox(Inches(1), Inches(3.95), Inches(7.6), Inches(0.3))
    tf = core_desc.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Orchestration • State Machine • Event Bus"
    run.font.size = Pt(12)
    run.font.color.rgb = RgbColor(220, 255, 240)
    run.font.name = "Arial"
    
    # Layer 3: Bottom modules
    modules_bottom = [
        ("🎯 Event Engine", "Rally detection, scoring, game flow"),
        ("📊 Analytics", "Live dashboards, statistics, cloud sync"),
    ]
    
    for i, (title, desc) in enumerate(modules_bottom):
        left = 1 + (i * 4)
        card = create_card(slide, left, 4.6, 3.6, 1, RgbColor(25, 40, 50))
        card.line.color.rgb = RgbColor(50, 80, 100)
        
        title_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(4.7), Inches(3.4), Inches(0.35))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.size = Pt(14)
        run.font.color.rgb = PRIMARY_LIGHT
        run.font.bold = True
        run.font.name = "Arial"
        
        desc_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(5.05), Inches(3.4), Inches(0.35))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(11)
        run.font.color.rgb = GRAY_500
        run.font.name = "Arial"


def create_differentiators_slide(prs):
    """Slide 8: Product - Differentiators"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_background(slide, DARK, DARK_SECONDARY)
    add_slide_header(slide, 8)
    add_slide_title(slide, "Why CourtNG is Different", 
                   "Not just another sports tech solution — a complete paradigm shift.",
                   accent_word="Different")
    
    diffs = [
        ("🚀", "Fully Autonomous", "Not an \"assist\" layer — CourtNG runs the entire match end-to-end, from first serve to game over."),
        ("🌐", "Sport-Agnostic", "Single platform configures for tennis, pickleball, badminton, futsal — and more to come."),
        ("⚡", "Edge-First Architecture", "On-premise inference with sub-50ms latency. Works offline. No cloud dependency."),
        ("🔧", "Modular & Upgradeable", "Vision, Rules, Events, and Analytics are swappable. AI models upgrade independently."),
    ]
    
    for i, (icon, title, desc) in enumerate(diffs):
        top = 2.3 + (i * 1.1)
        
        # Card
        card = create_card(slide, 0.5, top, 8.8, 0.95, RgbColor(25, 35, 50))
        card.line.color.rgb = RgbColor(50, 70, 100)
        
        # Icon with gradient background
        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), Inches(top + 0.2), Inches(0.55), Inches(0.55)
        )
        set_shape_gradient(icon_bg, PRIMARY, RgbColor(6, 182, 212))
        icon_bg.line.fill.background()
        icon_bg.adjustments[0] = 0.2
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(0.7), Inches(top + 0.22), Inches(0.55), Inches(0.55))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = icon
        run.font.size = Pt(20)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1.4), Inches(top + 0.15), Inches(7.5), Inches(0.35))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(15)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Arial"
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.4), Inches(top + 0.5), Inches(7.5), Inches(0.4))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(12)
        run.font.color.rgb = GRAY_400
        run.font.name = "Arial"


def create_business_model_slide(prs):
    """Slide 9: Business Model"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_background(slide, RgbColor(20, 20, 50), DARK_SECONDARY)
    add_slide_header(slide, 9)
    add_slide_title(slide, "Business Model", 
                   "SaaS subscription + hardware bundle = predictable revenue + high margins.")
    
    tiers = [
        ("STARTER", "$500", "/month per court", ["Single court", "Basic scoring", "Local dashboard", "Email support"]),
        ("GROWTH", "$1,000", "/month for 2-5 courts", ["2-5 courts", "Full analytics suite", "Mobile app access", "Leaderboards", "Priority support"]),
        ("ENTERPRISE", "$2,000+", "/month unlimited", ["Unlimited courts", "Custom rules engine", "Full API access", "Cloud sync", "Dedicated manager"]),
    ]
    
    for i, (tier, price, period, features) in enumerate(tiers):
        left = 0.5 + (i * 3.15)
        is_featured = (i == 1)
        
        # Card
        if is_featured:
            card = create_card(slide, left - 0.1, 2.1, 3.15, 3.2, PRIMARY)
            card.line.fill.background()
            tier_color = RgbColor(200, 255, 230)
            price_color = WHITE
            period_color = RgbColor(200, 255, 230)
            feature_color = RgbColor(220, 255, 240)
        else:
            card = create_card(slide, left, 2.3, 2.95, 2.95, RgbColor(30, 40, 60))
            card.line.color.rgb = RgbColor(60, 80, 120)
            tier_color = PRIMARY_LIGHT
            price_color = WHITE
            period_color = GRAY_500
            feature_color = GRAY_400
        
        top_offset = 2.1 if is_featured else 2.3
        card_left = left - 0.1 if is_featured else left
        
        # Tier name
        tier_box = slide.shapes.add_textbox(Inches(card_left + 0.2), Inches(top_offset + 0.2), Inches(2.7), Inches(0.3))
        tf = tier_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = tier
        run.font.size = Pt(11)
        run.font.color.rgb = tier_color
        run.font.bold = True
        run.font.name = "Arial"
        
        # Price
        price_box = slide.shapes.add_textbox(Inches(card_left + 0.2), Inches(top_offset + 0.5), Inches(2.7), Inches(0.5))
        tf = price_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = price
        run.font.size = Pt(32)
        run.font.color.rgb = price_color
        run.font.bold = True
        run.font.name = "Arial"
        
        # Period
        period_box = slide.shapes.add_textbox(Inches(card_left + 0.2), Inches(top_offset + 1), Inches(2.7), Inches(0.3))
        tf = period_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = period
        run.font.size = Pt(10)
        run.font.color.rgb = period_color
        run.font.name = "Arial"
        
        # Features
        feature_top = top_offset + 1.4
        for j, feat in enumerate(features):
            feat_box = slide.shapes.add_textbox(Inches(card_left + 0.2), Inches(feature_top + (j * 0.28)), Inches(2.7), Inches(0.28))
            tf = feat_box.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"✓ {feat}"
            run.font.size = Pt(10)
            run.font.color.rgb = feature_color
            run.font.name = "Arial"
    
    # Hardware note
    hw_bar = create_card(slide, 0.5, 5.45, 8.8, 0.45, RgbColor(30, 40, 60))
    hw_bar.line.color.rgb = RgbColor(60, 80, 120)
    
    hw_text = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(8.8), Inches(0.35))
    tf = hw_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Hardware Bundle: NVIDIA Jetson + Cameras + PoE Network = "
    run.font.size = Pt(12)
    run.font.color.rgb = GRAY_400
    run.font.name = "Arial"
    run2 = p.add_run()
    run2.text = "$8K–$15K"
    run2.font.size = Pt(12)
    run2.font.color.rgb = WHITE
    run2.font.bold = True
    run2.font.name = "Arial"
    run3 = p.add_run()
    run3.text = " one-time setup"
    run3.font.size = Pt(12)
    run3.font.color.rgb = GRAY_400
    run3.font.name = "Arial"


def create_competition_slide(prs):
    """Slide 10: Competition"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_background(slide, DARK, DARK_SECONDARY)
    add_slide_header(slide, 10)
    add_slide_title(slide, "Competitive Landscape", 
                   "Current solutions are expensive, limited, or both. We're changing that.")
    
    # Table header
    headers = ["Company", "Autonomous", "Multi-Sport", "Cost/Court", "Target"]
    header_widths = [2, 1.5, 1.5, 1.5, 1.8]
    
    # Header background
    header_bg = create_card(slide, 0.5, 2.3, 8.8, 0.5, RgbColor(35, 45, 65))
    header_bg.line.fill.background()
    
    left = 0.6
    for i, (header, width) in enumerate(zip(headers, header_widths)):
        hdr_box = slide.shapes.add_textbox(Inches(left), Inches(2.38), Inches(width), Inches(0.35))
        tf = hdr_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = header
        run.font.size = Pt(10)
        run.font.color.rgb = GRAY_400
        run.font.bold = True
        run.font.name = "Arial"
        left += width
    
    # Table rows
    competitors = [
        ("Hawk-Eye", "⚠️ Assist Only", "❌ 1-2 Sports", "$60K+", "Elite Only", RgbColor(239, 68, 68)),
        ("PlaySight", "⚠️ Assist Only", "❌ Tennis", "$20K", "Academies", RgbColor(245, 158, 11)),
        ("SwingVision", "⚠️ Assist Only", "❌ Tennis", "App", "Consumers", RgbColor(34, 197, 94)),
        ("Wingfield", "⚠️ Assist Only", "❌ Tennis", "$15K", "Clubs", RgbColor(245, 158, 11)),
    ]
    
    row_top = 2.85
    for company, auto, sport, cost, target, cost_color in competitors:
        left = 0.6
        row_data = [(company, WHITE, True), (auto, RgbColor(245, 158, 11), False), 
                    (sport, RgbColor(239, 68, 68), False), (cost, cost_color, True), (target, GRAY_400, False)]
        
        for j, ((text, color, bold), width) in enumerate(zip(row_data, header_widths)):
            cell = slide.shapes.add_textbox(Inches(left), Inches(row_top), Inches(width), Inches(0.4))
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(11)
            run.font.color.rgb = color
            run.font.bold = bold
            run.font.name = "Arial"
            left += width
        
        row_top += 0.5
    
    # CourtNG row (highlighted)
    courtng_bg = create_card(slide, 0.5, row_top - 0.05, 8.8, 0.55, RgbColor(20, 60, 50))
    courtng_bg.line.color.rgb = PRIMARY
    
    left = 0.6
    courtng_data = [("CourtNG", PRIMARY_LIGHT, True), ("✅ Full", PRIMARY, True), 
                    ("✅ 4+ Sports", PRIMARY, True), ("$8-15K", PRIMARY, True), ("Everyone", PRIMARY_LIGHT, True)]
    
    for (text, color, bold), width in zip(courtng_data, header_widths):
        cell = slide.shapes.add_textbox(Inches(left), Inches(row_top + 0.05), Inches(width), Inches(0.4))
        tf = cell.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(12)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Arial"
        left += width


def create_competitive_advantage_slide(prs):
    """Slide 11: Competitive Advantage"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_background(slide, RgbColor(15, 35, 35), DARK_SECONDARY)
    add_slide_header(slide, 11)
    add_slide_title(slide, "Our Competitive Moat", 
                   "Multiple defensible advantages that compound over time.",
                   accent_word="Moat")
    
    advantages = [
        ("🎯", "End-to-End Automation", "Hawk-Eye only does line calling. We own the entire match — from check-in to final score."),
        ("🌐", "Sport-Agnostic Platform", "Competitors are single-sport. Our rules engine scales across tennis, pickleball, badminton, futsal."),
        ("💰", "5-10x Cost Advantage", "Hawk-Eye: $60K/court. CourtNG: $8-15K. We unlock 99% of venues they can't serve."),
        ("🔒", "Full-Stack Integration", "No competitor automates officiating + scoring + UX together. This is our moat."),
    ]
    
    positions = [(0.5, 2.3), (5, 2.3), (0.5, 4.1), (5, 4.1)]
    
    for i, (icon, title, desc) in enumerate(advantages):
        left, top = positions[i]
        
        # Card
        card = create_card(slide, left, top, 4.3, 1.5, RgbColor(20, 50, 45))
        card.line.color.rgb = RgbColor(40, 100, 90)
        
        # Icon with gradient background
        icon_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + 0.2), Inches(top + 0.25), Inches(0.6), Inches(0.6)
        )
        set_shape_gradient(icon_bg, PRIMARY, RgbColor(6, 182, 212))
        icon_bg.line.fill.background()
        icon_bg.adjustments[0] = 0.25
        
        # Icon
        icon_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.28), Inches(0.6), Inches(0.6))
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = icon
        run.font.size = Pt(24)
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(left + 0.95), Inches(top + 0.2), Inches(3.2), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(15)
        run.font.color.rgb = WHITE
        run.font.bold = True
        run.font.name = "Arial"
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(left + 0.95), Inches(top + 0.6), Inches(3.2), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(11)
        run.font.color.rgb = GRAY_400
        run.font.name = "Arial"


def create_thank_you_slide(prs):
    """Slide 12: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_slide_background(slide, DARK, DARK_SECONDARY)
    add_slide_header(slide, 12)
    
    # Logo
    logo = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.15), Inches(1.3), Inches(1.5), Inches(1.5)
    )
    set_shape_gradient(logo, PRIMARY, RgbColor(6, 182, 212))
    logo.line.fill.background()
    logo.adjustments[0] = 0.3
    
    # Court icon
    icon_text = slide.shapes.add_textbox(Inches(4.15), Inches(1.6), Inches(1.5), Inches(0.9))
    tf = icon_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "⊞"
    run.font.size = Pt(48)
    run.font.color.rgb = WHITE
    
    # Thank You title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(0.9))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You"
    run.font.size = Pt(56)
    run.font.color.rgb = PRIMARY_LIGHT
    run.font.bold = True
    run.font.name = "Arial"
    
    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(3.85), Inches(9), Inches(0.4))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Let's build the future of sports together."
    run.font.size = Pt(20)
    run.font.color.rgb = GRAY_400
    run.font.name = "Arial"
    
    # Tagline
    tagline = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(0.4))
    tf = tagline.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "CourtNG — The Operating System for Courts That Run Themselves"
    run.font.size = Pt(14)
    run.font.color.rgb = GRAY_500
    run.font.name = "Arial"
    
    # Contact info
    contacts = [("📧 hello@courtng.com", 3.2), ("🌐 www.courtng.com", 5.7)]
    
    for text, left in contacts:
        contact_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left), Inches(4.9), Inches(2.3), Inches(0.45)
        )
        set_shape_solid(contact_bg, RgbColor(30, 40, 55))
        contact_bg.line.color.rgb = RgbColor(50, 60, 80)
        contact_bg.adjustments[0] = 0.5
        
        contact_text = slide.shapes.add_textbox(Inches(left), Inches(4.97), Inches(2.3), Inches(0.35))
        tf = contact_text.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(12)
        run.font.color.rgb = GRAY_400
        run.font.name = "Arial"
    
    # Funding badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(3.5), Inches(5.6), Inches(2.8), Inches(0.45)
    )
    set_shape_solid(badge, RgbColor(30, 45, 55))
    badge.line.color.rgb = RgbColor(50, 70, 90)
    badge.adjustments[0] = 0.5
    
    badge_text = slide.shapes.add_textbox(Inches(3.5), Inches(5.67), Inches(2.8), Inches(0.35))
    tf = badge_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "● Seeking $1M Seed Round"
    run.font.size = Pt(12)
    run.font.color.rgb = PRIMARY_LIGHT
    run.font.bold = True
    run.font.name = "Arial"


def main():
    """Create the full pitch deck"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(6.25)  # 16:10 aspect ratio
    
    print("Creating CourtNG Pitch Deck...")
    
    # Create all slides
    create_title_slide(prs)
    print("✓ Slide 1: Title")
    
    create_problem_slide(prs)
    print("✓ Slide 2: Problem")
    
    create_solution_slide(prs)
    print("✓ Slide 3: Solution")
    
    create_validation_slide(prs)
    print("✓ Slide 4: Market Validation")
    
    create_market_size_slide(prs)
    print("✓ Slide 5: Market Size")
    
    create_product_features_slide(prs)
    print("✓ Slide 6: Product Features")
    
    create_architecture_slide(prs)
    print("✓ Slide 7: Architecture")
    
    create_differentiators_slide(prs)
    print("✓ Slide 8: Differentiators")
    
    create_business_model_slide(prs)
    print("✓ Slide 9: Business Model")
    
    create_competition_slide(prs)
    print("✓ Slide 10: Competition")
    
    create_competitive_advantage_slide(prs)
    print("✓ Slide 11: Competitive Advantage")
    
    create_thank_you_slide(prs)
    print("✓ Slide 12: Thank You")
    
    # Save the presentation
    output_path = "CourtNG_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"\n✅ Pitch deck saved to: {output_path}")
    
    return output_path


if __name__ == "__main__":
    main()
